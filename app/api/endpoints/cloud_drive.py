import os
import uuid
import base64
import io
import logging
from typing import Optional, List
from datetime import datetime

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Query, Form
from pydantic import BaseModel
from sqlalchemy.orm import Session

from app.db.database import get_db
from app.api.dependencies import CurrentUser, get_current_user
from app.crud.cloud_file import cloud_file_crud, CLOUD_DIR
from app.models.cloud_file import CloudFile

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/cloud-drive", tags=["云盘"])

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB

ALLOWED_EXTENSIONS = {
    '.png', '.jpg', '.jpeg', '.gif', '.bmp',
    '.pdf', '.pptx', '.ppt', '.docx', '.doc',
    '.txt', '.md',
}

FILE_TYPE_MAP = {
    '.png': ('image', 'image/png'),
    '.jpg': ('image', 'image/jpeg'),
    '.jpeg': ('image', 'image/jpeg'),
    '.gif': ('image', 'image/gif'),
    '.bmp': ('image', 'image/bmp'),
    '.pdf': ('pdf', 'application/pdf'),
    '.pptx': ('pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'),
    '.ppt': ('ppt', 'application/vnd.ms-powerpoint'),
    '.docx': ('docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'),
    '.doc': ('doc', 'application/msword'),
    '.txt': ('txt', 'text/plain'),
    '.md': ('md', 'text/markdown'),
}

# 文件类型 → 上传 accept 扩展名映射
TYPE_ACCEPT_MAP = {
    '': '*',
    'docx': '.docx',
    'pptx': '.pptx',
    'ppt': '.ppt',
    'pdf': '.pdf',
    'image': '.png,.jpg,.jpeg,.gif,.bmp',
    'txt': '.txt,.md',
    'md': '.md',
    'doc': '.doc',
}


# ── Schemas ──

class CloudFileResponse(BaseModel):
    id: str
    file_name: str
    file_type: str
    mime_type: Optional[str] = None
    file_size: int
    is_folder: bool = False
    parent_id: Optional[str] = None
    created_at: str
    updated_at: str

    class Config:
        from_attributes = True


class CloudFileListResponse(BaseModel):
    files: List[CloudFileResponse]
    parent_id: Optional[str] = None
    folder_path: Optional[List[dict]] = None


class CloudFileDetailResponse(BaseModel):
    id: str
    file_name: str
    file_type: str
    mime_type: Optional[str] = None
    file_size: int
    is_folder: bool = False
    parent_id: Optional[str] = None
    content_text: Optional[str] = None
    base64: str
    created_at: str
    updated_at: str


class CreateFileRequest(BaseModel):
    file_name: str
    file_type: str  # pdf, docx, pptx
    parent_id: Optional[str] = None


class CreateFolderRequest(BaseModel):
    folder_name: str
    parent_id: Optional[str] = None


class CreateFromMarkdownRequest(BaseModel):
    file_name: str
    markdown: str
    parent_id: Optional[str] = None


class UpdateDocRequest(BaseModel):
    content_text: Optional[str] = None
    file_name: Optional[str] = None


class CreatePdfFromImagesRequest(BaseModel):
    file_name: str
    images: List[dict]  # [{base64, crop_x, crop_y, crop_w, crop_h, brightness, contrast}]
    parent_id: Optional[str] = None


# ── Helpers ──

def _to_response(f: CloudFile) -> CloudFileResponse:
    return CloudFileResponse(
        id=str(f.id),
        file_name=f.file_name,
        file_type=f.file_type,
        mime_type=f.mime_type,
        file_size=f.file_size or 0,
        is_folder=f.is_folder or False,
        parent_id=str(f.parent_id) if f.parent_id else None,
        created_at=f.created_at.isoformat() if f.created_at else '',
        updated_at=f.updated_at.isoformat() if f.updated_at else '',
    )


def _create_empty_docx() -> bytes:
    from docx import Document
    doc = Document()
    doc.add_paragraph('新建文档', style='Title')
    doc.add_paragraph('')
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _create_empty_pdf() -> bytes:
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=16)
    pdf.cell(0, 20, text="新建文档", new_x="LMARGIN", new_y="NEXT", align="C")
    buf = io.BytesIO()
    pdf.output(buf)
    return buf.getvalue()


def _markdown_to_docx(markdown: str) -> bytes:
    """将 Markdown 转为 DOCX"""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    import re

    doc = Document()

    # Set default font
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Microsoft YaHei'
    font.size = Pt(11)

    lines = markdown.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]

        # 标题
        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[2:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[2:], level=3)
        elif line.startswith('---'):
            doc.add_paragraph('').add_run('─' * 40)
        elif line.strip() == '':
            doc.add_paragraph('')
        elif line.startswith('- ') or line.startswith('* '):
            # 简单列表
            p = doc.add_paragraph(style='List Bullet')
            p.add_run(line[2:])
        elif re.match(r'^\d+[\.)] ', line):
            p = doc.add_paragraph(style='List Number')
            p.add_run(re.sub(r'^\d+[\.)] ', '', line))
        else:
            # 普通段落，支持加粗和斜体
            p = doc.add_paragraph()
            # 简单处理加粗 **text**
            parts = re.split(r'(\*\*.*?\*\*)', line)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run(part[2:-2])
                    run.bold = True
                else:
                    p.add_run(part)

        i += 1

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── Endpoints ──

@router.get("/files", response_model=CloudFileListResponse)
async def list_cloud_files(
    parent_id: Optional[str] = Query(None, description="文件夹ID，不传则列出根目录"),
    file_type: Optional[str] = Query(None, description="按文件类型过滤"),
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """列出文件夹内容"""
    files = cloud_file_crud.list_files(db, str(current_user.student_id), parent_id=parent_id, file_type=file_type)
    folder_path = None
    if parent_id:
        folder_path = cloud_file_crud.get_folder_path(db, parent_id, str(current_user.student_id))
    return CloudFileListResponse(
        files=[_to_response(f) for f in files],
        parent_id=parent_id,
        folder_path=folder_path,
    )


@router.post("/folder")
async def create_folder(
    req: CreateFolderRequest,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """新建文件夹"""
    if not req.folder_name.strip():
        raise HTTPException(status_code=400, detail="文件夹名称不能为空")

    # 检查同目录下是否重名
    existing = cloud_file_crud.list_files(db, str(current_user.student_id), parent_id=req.parent_id)
    for f in existing:
        if f.is_folder and f.file_name == req.folder_name.strip():
            raise HTTPException(status_code=400, detail="该文件夹名称已存在")

    record = cloud_file_crud.create_file(
        db,
        user_id=str(current_user.student_id),
        file_name=req.folder_name.strip(),
        file_type='folder',
        mime_type=None,
        file_size=0,
        storage_key=None,
        parent_id=req.parent_id,
        is_folder=True,
    )
    return _to_response(record)


@router.post("/upload", response_model=CloudFileResponse)
async def upload_cloud_file(
    file: UploadFile = File(...),
    parent_id: Optional[str] = Form(None),
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """上传文件到云盘"""
    if not file.filename:
        raise HTTPException(status_code=400, detail="文件名不能为空")

    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件类型 ({ext})，仅支持: {', '.join(sorted(ALLOWED_EXTENSIONS))}",
        )

    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="文件过大，单个文件不能超过 50MB")

    file_type, mime_type = FILE_TYPE_MAP.get(ext, ('unknown', 'application/octet-stream'))
    storage_key = f"{uuid.uuid4().hex}{ext}"
    filepath = os.path.join(CLOUD_DIR, storage_key)

    os.makedirs(CLOUD_DIR, exist_ok=True)
    with open(filepath, 'wb') as f:
        f.write(content)

    record = cloud_file_crud.create_file(
        db,
        user_id=str(current_user.student_id),
        file_name=file.filename,
        file_type=file_type,
        mime_type=mime_type,
        file_size=len(content),
        storage_key=storage_key,
        parent_id=parent_id,
    )
    return _to_response(record)


@router.post("/create", response_model=CloudFileResponse)
async def create_cloud_file(
    req: CreateFileRequest,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """新建空白文档"""
    file_type = req.file_type.lower()
    if file_type not in ('docx', 'pdf'):
        raise HTTPException(status_code=400, detail="仅支持新建 docx、pdf 文件")

    if not req.file_name.strip():
        raise HTTPException(status_code=400, detail="文件名不能为空")

    base_name = req.file_name.rsplit('.', 1)[0] if '.' in req.file_name else req.file_name
    full_name = f"{base_name}.{file_type}"

    if file_type == 'docx':
        content = _create_empty_docx()
        mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    else:
        content = _create_empty_pdf()
        mime_type = 'application/pdf'

    storage_key = f"{uuid.uuid4().hex}.{file_type}"
    filepath = os.path.join(CLOUD_DIR, storage_key)
    os.makedirs(CLOUD_DIR, exist_ok=True)
    with open(filepath, 'wb') as f:
        f.write(content)

    record = cloud_file_crud.create_file(
        db,
        user_id=str(current_user.student_id),
        file_name=full_name,
        file_type=file_type,
        mime_type=mime_type,
        file_size=len(content),
        storage_key=storage_key,
        parent_id=req.parent_id,
    )
    return _to_response(record)


@router.post("/create-from-markdown", response_model=CloudFileResponse)
async def create_word_from_markdown(
    req: CreateFromMarkdownRequest,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """从 Markdown 内容生成 Word 文档"""
    if not req.file_name.strip():
        raise HTTPException(status_code=400, detail="文件名不能为空")
    if not req.markdown.strip():
        raise HTTPException(status_code=400, detail="内容不能为空")

    base_name = req.file_name.rsplit('.', 1)[0] if '.' in req.file_name else req.file_name
    full_name = f"{base_name}.docx"

    content = _markdown_to_docx(req.markdown)
    storage_key = f"{uuid.uuid4().hex}.docx"
    filepath = os.path.join(CLOUD_DIR, storage_key)
    os.makedirs(CLOUD_DIR, exist_ok=True)
    with open(filepath, 'wb') as f:
        f.write(content)

    record = cloud_file_crud.create_file(
        db,
        user_id=str(current_user.student_id),
        file_name=full_name,
        file_type='docx',
        mime_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        file_size=len(content),
        storage_key=storage_key,
        parent_id=req.parent_id,
        content_text=req.markdown,
    )
    return _to_response(record)


@router.put("/files/{file_id}", response_model=CloudFileResponse)
async def update_cloud_file(
    file_id: str,
    req: UpdateDocRequest,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """更新文档内容（Word 编辑保存）"""
    record = cloud_file_crud.get_file(db, file_id, str(current_user.student_id))
    if not record:
        raise HTTPException(status_code=404, detail="文件不存在")

    # 如果是 docx 且有新的 markdown 内容，重新生成 docx 文件
    if req.content_text is not None and record.file_type == 'docx':
        content = _markdown_to_docx(req.content_text)
        if record.storage_key:
            filepath = cloud_file_crud.get_storage_path(record.storage_key)
            with open(filepath, 'wb') as f:
                f.write(content)
            record.file_size = len(content)

    updated = cloud_file_crud.update_file_content(
        db, file_id, str(current_user.student_id),
        content_text=req.content_text,
        file_name=req.file_name,
    )
    if not updated:
        raise HTTPException(status_code=404, detail="文件不存在")
    return _to_response(updated)


@router.post("/create-pdf-from-images", response_model=CloudFileResponse)
async def create_pdf_from_images(
    req: CreatePdfFromImagesRequest,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """从图片生成 PDF（支持裁切、亮度、对比度调整）"""
    if not req.file_name.strip():
        raise HTTPException(status_code=400, detail="文件名不能为空")
    if not req.images:
        raise HTTPException(status_code=400, detail="请至少添加一张图片")

    from PIL import Image, ImageEnhance, ImageFilter
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=False)

    for img_data in req.images:
        b64 = img_data.get("base64", "")
        if not b64:
            continue

        try:
            img_bytes = base64.b64decode(b64)
            img = Image.open(io.BytesIO(img_bytes))

            # 裁切
            crop_x = img_data.get("crop_x")
            crop_y = img_data.get("crop_y")
            crop_w = img_data.get("crop_w")
            crop_h = img_data.get("crop_h")
            if all(v is not None for v in [crop_x, crop_y, crop_w, crop_h]):
                img = img.crop((crop_x, crop_y, crop_x + crop_w, crop_y + crop_h))

            # 亮度
            brightness = img_data.get("brightness", 1.0)
            if brightness != 1.0:
                img = ImageEnhance.Brightness(img).enhance(brightness)

            # 对比度
            contrast = img_data.get("contrast", 1.0)
            if contrast != 1.0:
                img = ImageEnhance.Contrast(img).enhance(contrast)

            # 转为 RGB 并保存为临时文件
            if img.mode != 'RGB':
                img = img.convert('RGB')

            temp_path = os.path.join(CLOUD_DIR, f"_temp_{uuid.uuid4().hex}.jpg")
            img.save(temp_path, 'JPEG', quality=85)

            # 计算适合页面的尺寸（A4: 210x297mm）
            pdf.add_page()
            max_w = 190  # mm
            max_h = 270  # mm
            img_w, img_h = img.size
            ratio = min(max_w / img_w, max_h / img_h)
            pdf_w = img_w * ratio
            pdf_h = img_h * ratio
            x = (210 - pdf_w) / 2
            y = (297 - pdf_h) / 2
            pdf.image(temp_path, x=x, y=y, w=pdf_w, h=pdf_h)

            if os.path.exists(temp_path):
                os.remove(temp_path)

        except Exception as e:
            logger.error(f"图片处理失败: {e}")
            continue

    storage_key = f"{uuid.uuid4().hex}.pdf"
    filepath = os.path.join(CLOUD_DIR, storage_key)
    os.makedirs(CLOUD_DIR, exist_ok=True)
    pdf.output(filepath)
    file_size = os.path.getsize(filepath)

    record = cloud_file_crud.create_file(
        db,
        user_id=str(current_user.student_id),
        file_name=req.file_name.strip() + '.pdf' if '.' not in req.file_name else req.file_name,
        file_type='pdf',
        mime_type='application/pdf',
        file_size=file_size,
        storage_key=storage_key,
        parent_id=req.parent_id,
    )
    return _to_response(record)


@router.get("/files/{file_id}", response_model=CloudFileDetailResponse)
async def get_cloud_file(
    file_id: str,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """获取文件详情"""
    record = cloud_file_crud.get_file(db, file_id, str(current_user.student_id))
    if not record:
        raise HTTPException(status_code=404, detail="文件不存在")

    content = b''
    if record.storage_key:
        filepath = cloud_file_crud.get_storage_path(record.storage_key)
        if os.path.exists(filepath):
            with open(filepath, 'rb') as f:
                content = f.read()

    return CloudFileDetailResponse(
        id=str(record.id),
        file_name=record.file_name,
        file_type=record.file_type,
        mime_type=record.mime_type,
        file_size=record.file_size or 0,
        is_folder=record.is_folder or False,
        parent_id=str(record.parent_id) if record.parent_id else None,
        content_text=record.content_text,
        base64=base64.b64encode(content).decode('utf-8') if content else '',
        created_at=record.created_at.isoformat() if record.created_at else '',
        updated_at=record.updated_at.isoformat() if record.updated_at else '',
    )


@router.get("/files/{file_id}/download")
async def download_cloud_file(
    file_id: str,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """下载文件"""
    record = cloud_file_crud.get_file(db, file_id, str(current_user.student_id))
    if not record:
        raise HTTPException(status_code=404, detail="文件不存在")

    if record.is_folder:
        raise HTTPException(status_code=400, detail="文件夹不能下载")

    filepath = cloud_file_crud.get_storage_path(record.storage_key)
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="文件数据不存在")

    from fastapi.responses import FileResponse
    return FileResponse(filepath, filename=record.file_name)


@router.get("/files/{file_id}/preview-text")
async def preview_cloud_file_text(
    file_id: str,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """获取文件预览文本"""
    record = cloud_file_crud.get_file(db, file_id, str(current_user.student_id))
    if not record:
        raise HTTPException(status_code=404, detail="文件不存在")

    result = {"text": "", "slides": None, "content_text": record.content_text}

    # Word 文档优先返回保存的 markdown 源
    if record.file_type == 'docx' and record.content_text:
        result["text"] = record.content_text
        return result

    filepath = cloud_file_crud.get_storage_path(record.storage_key) if record.storage_key else None
    if not filepath or not os.path.exists(filepath):
        return result

    if record.file_type == 'docx':
        from app.services.docx_parser import extract_docx_text
        result["text"] = extract_docx_text(filepath, max_chars=0)

    elif record.file_type in ("pptx", "ppt"):
        from pptx import Presentation
        prs = Presentation(filepath)
        slides_text = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
            slides_text.append(slide_texts)
        result["slides"] = slides_text

    elif record.file_type == "pdf":
        import fitz
        doc = fitz.open(filepath)
        pages = []
        for page_num in range(len(doc)):
            text = doc[page_num].get_text().strip()
            if text:
                pages.append(text)
        doc.close()
        result["text"] = "\n\n".join(pages)

    else:
        try:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                result["text"] = f.read()
        except:
            result["text"] = ""

    return result


@router.delete("/files/{file_id}")
async def delete_cloud_file(
    file_id: str,
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """删除云盘文件"""
    success = cloud_file_crud.delete_file(db, file_id, str(current_user.student_id))
    if not success:
        raise HTTPException(status_code=404, detail="文件不存在")
    return {"success": True}


@router.get("/files/{file_id}/thumbnail")
async def get_file_thumbnail(
    file_id: str,
    width: int = Query(200, description="缩略图宽度（仅图片有效）"),
    current_user: CurrentUser = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    """获取文件缩略图：图片返回缩小后的图片，PPT返回第一页文字"""
    record = cloud_file_crud.get_file(db, file_id, str(current_user.student_id))
    if not record:
        raise HTTPException(status_code=404, detail="文件不存在")

    if record.file_type == 'image' and record.storage_key:
        filepath = cloud_file_crud.get_storage_path(record.storage_key)
        if not os.path.exists(filepath):
            raise HTTPException(status_code=404, detail="文件数据不存在")
        try:
            from PIL import Image
            import io
            img = Image.open(filepath)
            # 转换 RGBA → RGB（JPEG 不支持透明通道）
            if img.mode == 'RGBA':
                background = Image.new('RGB', img.size, (255, 255, 255))
                background.paste(img, mask=img.split()[3])
                img = background
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            # 等比缩放
            w_percent = width / float(img.size[0])
            height = int(float(img.size[1]) * w_percent)
            img = img.resize((width, height), Image.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format='JPEG', quality=70)
            from fastapi.responses import Response
            return Response(content=buf.getvalue(), media_type='image/jpeg')
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"生成缩略图失败: {str(e)}")

    if record.file_type in ('pptx', 'ppt') and record.storage_key:
        filepath = cloud_file_crud.get_storage_path(record.storage_key)
        if not os.path.exists(filepath):
            raise HTTPException(status_code=404, detail="文件数据不存在")
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            first_slide = None
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                if texts:
                    first_slide = texts[:5]  # 最多5条
                break
            return {"type": "slide_text", "texts": first_slide or ["(空白幻灯片)"]}
        except Exception as e:
            return {"type": "slide_text", "texts": ["(无法解析)"]}

    return {"type": "none"}


@router.get("/accept-map")
async def get_accept_map():
    """获取文件类型对应的 accept 值"""
    return TYPE_ACCEPT_MAP
